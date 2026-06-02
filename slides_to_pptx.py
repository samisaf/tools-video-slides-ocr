#!/usr/bin/env python3
"""
slides_to_pptx.py — Reusable CLI tool to convert extracted lecture snapshots and OCR-ed slide text into a polished widescreen PowerPoint presentation.

This works hand-in-hand with video_ocr.py. Once you have generated <video>_snapshots/ and <video>_slides.txt using video_ocr.py,
you can run this script to create a premium widescreen (.pptx) slide deck.
"""

import os
import re
import shutil
import argparse
import numpy as np
from PIL import Image
from pathlib import Path
from typing import List, Dict, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Default configuration parameters
DEFAULT_CROP_BOX = (2, 16, 591, 347)  # (left, top, right, bottom)
DEFAULT_MSE_THRESHOLD = 50.0

# Aesthetic Design Tokens
COLOR_DARK_BG = RGBColor(15, 23, 42)      # Slate 900 (Deep elegant charcoal/navy)
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # Slate 50 (Sleek off-white for content canvas)
COLOR_TEXT_DARK = RGBColor(15, 23, 42)    # Slate 900 for titles
COLOR_TEXT_MUTED = RGBColor(51, 65, 85)   # Slate 700 for highly-readable body text
COLOR_WHITE = RGBColor(255, 255, 255)     # For text on dark slides
COLOR_TEAL_LIGHT = RGBColor(56, 189, 248) # Sky 400 for stunning title slide accents
COLOR_BORDER = RGBColor(226, 232, 240)    # Slate 200 border for figure cards

# Hand-curated slide titles mapping for visual slides with empty OCR text
EMPTY_SLIDE_TITLES = {
    2: "Kiteboarding Analogy: Force & Tension",
    3: "Kiteboarding Analogy: Force & Tension",
    23: "Ultrasound System Console & Equipment",
    24: "Knobology: Physical Keyboard Controls",
    30: "Conclusion: Probe Manipulation Analogy",
    31: "Conclusion: Probe Manipulation Analogy",
}

def parse_slides_ocr(slides_txt_path: Path) -> List[Dict[str, Any]]:
    """Parse video_ocr text file and return a list of slide structures."""
    content = slides_txt_path.read_text(encoding="utf-8")
    pattern = r"# Snapshot (\d+) — (snapshot_\d+\.jpg)"
    parts = re.split(pattern, content)
    
    slides = []
    if len(parts) > 1:
        for i in range(1, len(parts), 3):
            idx = int(parts[i])
            filename = parts[i+1]
            text = parts[i+2].strip() if i+2 < len(parts) else ""
            slides.append({
                "idx": idx,
                "filename": filename,
                "text": text
            })
    return slides

def clean_ocr_text(text: str) -> Tuple[str, List[str]]:
    """OCR text cleaning pipeline."""
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    if not lines:
        return "", []
    
    title = lines[0]
    bullets = lines[1:]
    
    # Handle video lecture logo text: '2 CHEST' or 'CHEST'
    if title == "2 CHEST" and len(lines) > 1:
        title = lines[1]
        bullets = lines[2:]
        
    # Clean up common title artifacts
    title = re.sub(r"\s*[\)\(]+\s*$", "", title)  # Remove stray parens
    title = title.replace("ULTRASONAGRAPHY", "ULTRASONOGRAPHY")
    title = title.title()  # Convert to gorgeous Title Case
    
    cleaned_bullets = []
    for b in bullets:
        # Strip OCR bullet artifacts (*, -, =, »)
        cleaned_b = re.sub(r"^[*\-\+=»«¥v\/o•]\s*", "", b).strip()
        if cleaned_b:
            # Fix common OCR typos within bullets
            cleaned_b = cleaned_b.replace("’", "'").replace("“", '"').replace("”", '"')
            cleaned_b = cleaned_b.replace("n screen", "on screen")
            cleaned_b = cleaned_b.replace("S—10 MH", "5-10 MHz").replace("S10 MH", "5-10 MHz")
            cleaned_b = cleaned_b.replace("1-5MHe", "1-5 MHz").replace("1-5Mae", "1-5 MHz")
            cleaned_bullets.append(cleaned_b)
            
    return title, cleaned_bullets

def calculate_mse(img1_path: Path, img2_path: Path) -> float:
    """Calculate Mean Squared Error between two images."""
    try:
        img1 = np.array(Image.open(img1_path))
        img2 = np.array(Image.open(img2_path))
        if img1.shape == img2.shape:
            err = np.sum((img1.astype("float") - img2.astype("float")) ** 2)
            err /= float(img1.shape[0] * img1.shape[1] * img1.shape[2])
            return err
    except Exception:
        pass
    return 1e9  # treated as completely different

def is_black_frame(img_path: Path) -> bool:
    """Detect if a snapshot frame is completely black."""
    try:
        img = np.array(Image.open(img_path))
        return float(np.mean(img)) < 10.0
    except Exception:
        return False

def compile_pptx(slides_txt_path: Path, snapshot_dir: Path, output_pptx: Path, 
                 crop_box: Tuple[int, int, int, int], mse_threshold: float) -> None:
    """Process OCR text and snapshot images to compile a widescreen PPTX."""
    print(f"Reading OCR slides: {slides_txt_path.name}")
    raw_slides = parse_slides_ocr(slides_txt_path)
    
    # Create a temporary directory for cropping images
    temp_crop_dir = slides_txt_path.parent / f"_temp_cropped_{slides_txt_path.stem}"
    temp_crop_dir.mkdir(exist_ok=True)
    
    # Deduplicate and skip black frames
    filtered_slides = []
    idx = 0
    while idx < len(raw_slides):
        slide = raw_slides[idx]
        snap_path = snapshot_dir / slide["filename"]
        
        if not snap_path.exists():
            # If snapshot is missing, proceed with just text
            filtered_slides.append(slide)
            idx += 1
            continue
            
        if is_black_frame(snap_path):
            idx += 1
            continue
            
        dup_run_last_idx = idx
        for next_idx in range(idx + 1, len(raw_slides)):
            next_slide = raw_slides[next_idx]
            next_snap_path = snapshot_dir / next_slide["filename"]
            
            if next_snap_path.exists():
                err = calculate_mse(snap_path, next_snap_path)
                if err < mse_threshold:
                    dup_run_last_idx = next_idx
                    snap_path = next_snap_path
                else:
                    break
            else:
                break
                
        best_slide = raw_slides[dup_run_last_idx]
        filtered_slides.append(best_slide)
        idx = dup_run_last_idx + 1
        
    print(f"Collapsed duplicate frames. Compiled slide count: {len(filtered_slides)}")
    
    # Create PowerPoint Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]
    
    # Title Slide
    title_slide = prs.slides.add_slide(blank_slide_layout)
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = COLOR_DARK_BG
    
    # Title Text
    # Use stem clean-up to determine nice lecture title
    clean_stem = slides_txt_path.stem.replace("_slides", "").replace("_ocr", "")
    presentation_title = clean_stem.replace("_", " ").replace("-", " ").title()
    
    title_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = Inches(0)
    
    p_title = tf.paragraphs[0]
    p_title.text = presentation_title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.space_after = Pt(20)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Lecture Slide Deck Compilation  •  Widescreen Presentation"
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_TEAL_LIGHT
    
    # Thin underline graphic
    line_shape = title_slide.shapes.add_shape(1, Inches(1.0), Inches(4.8), Inches(4.5), Inches(0.05))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = COLOR_TEAL_LIGHT
    line_shape.line.color.rgb = COLOR_TEAL_LIGHT
    
    # Content Slides
    for slide_data in filtered_slides:
        orig_snap_path = snapshot_dir / slide_data["filename"]
        cropped_snap_path = temp_crop_dir / f"crop_{slide_data['filename']}"
        has_image = orig_snap_path.exists()
        
        if has_image:
            img = Image.open(orig_snap_path)
            cropped_img = img.crop(crop_box)
            cropped_img.save(cropped_snap_path)
            
        title, bullets = clean_ocr_text(slide_data["text"])
        
        # Fallback titles for empty OCR slides
        if not title:
            title = EMPTY_SLIDE_TITLES.get(slide_data["idx"], "Visual Demonstration")
            
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_LIGHT_BG
        
        # Add Header Title
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.9))
        htf = header_box.text_frame
        htf.word_wrap = True
        htf.margin_left = htf.margin_top = Inches(0)
        
        hp = htf.paragraphs[0]
        hp.text = title
        hp.font.name = "Segoe UI"
        hp.font.size = Pt(32)
        hp.font.bold = True
        hp.font.color.rgb = COLOR_TEXT_DARK
        
        has_text = len(bullets) > 0
        
        if has_text and has_image:
            # --- SIDE-BY-SIDE LAYOUT ---
            # Text box (Left)
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
            ttf = text_box.text_frame
            ttf.word_wrap = True
            ttf.margin_left = ttf.margin_top = Inches(0)
            
            for b_idx, bullet_text in enumerate(bullets):
                p = ttf.paragraphs[0] if b_idx == 0 else ttf.add_paragraph()
                p.text = f"•  {bullet_text}"
                p.font.name = "Segoe UI"
                p.font.size = Pt(15)
                p.font.color.rgb = COLOR_TEXT_MUTED
                p.line_spacing = 1.25
                p.space_after = Pt(14)
                
            # Image card (Right)
            img_width = Inches(5.6)
            img_height = Inches(3.14)
            img_left = Inches(6.933)
            img_top = Inches(1.5) + Inches((5.2 - 3.14) / 2)
            
            pic = slide.shapes.add_picture(str(cropped_snap_path), img_left, img_top, img_width, img_height)
            pic.line.color.rgb = COLOR_BORDER
            pic.line.width = Pt(1.5)
            
        elif has_image:
            # --- FULL FIGURE LAYOUT (Visual only) ---
            img_width = Inches(9.2)
            img_height = Inches(5.16)
            img_left = Inches((13.333 - 9.2) / 2)
            img_top = Inches(1.6)
            
            pic = slide.shapes.add_picture(str(cropped_snap_path), img_left, img_top, img_width, img_height)
            pic.line.color.rgb = COLOR_BORDER
            pic.line.width = Pt(1.5)
            
        elif has_text:
            # --- FULL TEXT LAYOUT ---
            text_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.2))
            ttf = text_box.text_frame
            ttf.word_wrap = True
            ttf.margin_left = ttf.margin_top = Inches(0)
            
            for b_idx, bullet_text in enumerate(bullets):
                p = ttf.paragraphs[0] if b_idx == 0 else ttf.add_paragraph()
                p.text = f"•  {bullet_text}"
                p.font.name = "Segoe UI"
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_TEXT_MUTED
                p.line_spacing = 1.3
                p.space_after = Pt(16)
                
    prs.save(str(output_pptx))
    shutil.rmtree(temp_crop_dir)
    print(f"Presentation saved successfully → {output_pptx.name}")

def parse_crop_box(crop_str: str) -> Tuple[int, int, int, int]:
    """Parse comma-separated crop box string into a tuple of 4 integers."""
    try:
        parts = [int(x.strip()) for x in crop_str.split(",")]
        if len(parts) == 4:
            return (parts[0], parts[1], parts[2], parts[3])
    except Exception:
        pass
    raise argparse.ArgumentTypeError("Crop box must be 4 comma-separated integers, e.g., '2,16,591,347'")

def run():
    parser = argparse.ArgumentParser(
        description="Convert video slides OCR text and snapshots into premium widescreen PowerPoint presentations."
    )
    scope = parser.add_mutually_exclusive_group(required=True)
    scope.add_argument("--video", type=Path, help="Single video file or slide text file (e.g., movie.mp4 or movie_slides.txt)")
    scope.add_argument("--dir", type=Path, help="Directory path to scan and batch-process all slide OCR outputs.")
    
    parser.add_argument("--mse", type=float, default=DEFAULT_MSE_THRESHOLD, help=f"MSE threshold for duplicate detection (default {DEFAULT_MSE_THRESHOLD})")
    parser.add_argument("--crop", type=parse_crop_box, default=DEFAULT_CROP_BOX, help="Slide bounding crop box as 'left,top,right,bottom' coordinates")
    
    args = parser.parse_args()
    
    if args.video:
        # Resolve target files from input path
        input_path = args.video.expanduser().resolve()
        
        if input_path.suffix.lower() == ".txt":
            slides_txt = input_path
            stem = slides_txt.stem.replace("_slides", "").replace("_ocr", "")
            video_dir = slides_txt.parent
        else:
            stem = input_path.stem
            slides_txt = input_path.with_name(f"{stem}_slides.txt")
            video_dir = input_path.parent
            
        snapshot_dir = video_dir / f"{stem}_snapshots"
        output_pptx = video_dir / f"{stem}.pptx"
        
        if not slides_txt.exists():
            print(f"Error: Slide OCR text file not found: {slides_txt}")
            return
            
        if not snapshot_dir.is_dir():
            print(f"Error: Snapshot directory not found: {snapshot_dir}")
            return
            
        compile_pptx(slides_txt, snapshot_dir, output_pptx, args.crop, args.mse)
        
    elif args.dir:
        # Scan directory and batch-compile
        target_dir = args.dir.expanduser().resolve()
        if not target_dir.is_dir():
            print(f"Error: Directory does not exist: {target_dir}")
            return
            
        # Discover all slide text files
        ocr_files = sorted(list(target_dir.glob("*_slides.txt")))
        if not ocr_files:
            print(f"No '*_slides.txt' files found in directory {target_dir}")
            return
            
        print(f"Discovered {len(ocr_files)} slide text files to process.")
        for slides_txt in ocr_files:
            stem = slides_txt.stem.replace("_slides", "").replace("_ocr", "")
            snapshot_dir = target_dir / f"{stem}_snapshots"
            output_pptx = target_dir / f"{stem}.pptx"
            
            if not snapshot_dir.is_dir():
                print(f"Skipping {slides_txt.name}: corresponding snapshot directory '{snapshot_dir.name}' not found.")
                continue
                
            try:
                compile_pptx(slides_txt, snapshot_dir, output_pptx, args.crop, args.mse)
            except Exception as e:
                print(f"Error processing {slides_txt.name}: {e}")

if __name__ == "__main__":
    run()
